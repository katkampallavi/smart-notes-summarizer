import os
import re


def extract_text_from_pptx(filepath):
    """
    Extract text content from a .pptx file using python-pptx.
    Extracts text from all slides, shapes, tables, and notes.
    Returns extracted text as a string.
    """
    if not os.path.exists(filepath):
        raise FileNotFoundError(f"PPTX file not found: {filepath}")

    try:
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation(filepath)
        all_slides_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            slide_texts.append(f"[Slide {slide_num}]")

            for shape in slide.shapes:
                # Extract from text frames
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = ''.join(run.text for run in paragraph.runs).strip()
                        if para_text:
                            slide_texts.append(para_text)

                # Extract from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_texts = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_texts.append(cell_text)
                        if row_texts:
                            slide_texts.append(' | '.join(row_texts))

            # Extract speaker notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame:
                    notes_text = notes_frame.text.strip()
                    if notes_text:
                        slide_texts.append(f"(Notes: {notes_text})")

            if len(slide_texts) > 1:  # More than just the slide label
                all_slides_text.append('\n'.join(slide_texts))

        text = '\n\n'.join(all_slides_text)

        if not text.strip():
            raise ValueError("The PPTX file appears to be empty or contains no extractable text.")

        return _clean_extracted_text(text)

    except ImportError:
        raise ImportError(
            "python-pptx is required to read PPTX files. "
            "Install it with: pip install python-pptx"
        )
    except Exception as e:
        raise RuntimeError(f"Error reading PPTX file: {str(e)}")


def _clean_extracted_text(text):
    """Clean and normalize extracted text from PPTX."""
    # Remove non-printable characters except common whitespace
    text = re.sub(r'[^\x20-\x7E\n\r\t]', ' ', text)
    # Normalize spaces
    text = re.sub(r'[ \t]+', ' ', text)
    # Normalize newlines
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = text.strip()
    return text